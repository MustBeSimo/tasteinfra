from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE

# --- CONFIGURATION ---
ACCENT_COLOR = RGBColor(245, 158, 11)  # #f59e0b
DARK_BG = RGBColor(26, 26, 26)         # #1a1a1a
LIGHT_BG = RGBColor(255, 255, 255)     # #ffffff
TEXT_DARK = RGBColor(255, 255, 255)
TEXT_LIGHT = RGBColor(0, 0, 0)
TEXT_GRAY = RGBColor(128, 128, 128)

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, text, left, top, width, height, font_size=18, color=TEXT_LIGHT, bold=False, align=None, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    if align:
        p.alignment = align
    return txBox

def clean_text(text):
    return text.replace("\n", " ").strip()

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333) # 16:9 Aspect Ratio
prs.slide_height = Inches(7.5)

# --------------------------------------------------------------------------------
# SLIDE 1: TITLE (Dark)
# --------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
set_slide_background(slide, DARK_BG)

# Label
add_text_box(slide, "Manifesto — Studio W230", Inches(1), Inches(1), Inches(10), Inches(0.5), 14, TEXT_GRAY)

# Title
add_text_box(slide, "The AI Revolution is Over.\nThe Evaluation Revolution Just Started.", 
             Inches(1), Inches(1.8), Inches(11), Inches(2), 44, TEXT_DARK, bold=True)

# Subtitle
sub = add_text_box(slide, "When generation gets cheap, evaluation becomes the edge.", 
             Inches(1), Inches(4), Inches(11), Inches(1), 24, TEXT_DARK)
sub.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255) 
# Note: Complex rich text formatting (orange span) needs manual run manipulation, simplified here for script brevity.

# Footer info
add_text_box(slide, "Simone Leonelli — 16 years building taste systems for luxury brands.", 
             Inches(1), Inches(5.5), Inches(8), Inches(1), 14, TEXT_GRAY)
add_text_box(slide, "January 2026   v3.2", Inches(1), Inches(6.5), Inches(4), Inches(0.5), 14, TEXT_GRAY)


# --------------------------------------------------------------------------------
# SLIDE 2: THESIS (Light)
# --------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)

add_text_box(slide, "01 — The Warning", Inches(1), Inches(0.5), Inches(5), Inches(0.5), 12, TEXT_GRAY)
add_text_box(slide, "The \"Generation at Scale\" Lie.", Inches(1), Inches(1), Inches(10), Inches(1), 36, TEXT_LIGHT, bold=True)

# Callout Box (Dark Box on Light Slide)
shape = slide.shapes.add_shape(1, Inches(1), Inches(2.2), Inches(11.3), Inches(2)) # Rectangle
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BG
shape.line.color.rgb = ACCENT_COLOR
shape.line.width = Pt(3)
tf = shape.text_frame
p = tf.paragraphs[0]
p.text = "Q3 2025: The Corpse\nA Fortune 500 retailer shipped 50,000 AI-generated product emails. 12% contained hallucinations. Revenue dropped 15%.\nThis is the Taste Infrastructure problem."
p.font.color.rgb = RGBColor(200, 200, 200)
p.font.size = Pt(16)

# Two Col
add_text_box(slide, "The Villain\nVibe-Driven Deployment\n\"It looks fine to me.\"", 
             Inches(1), Inches(5), Inches(5), Inches(2), 18, TEXT_LIGHT)

# Hero Box (Orange accent)
hero_box = slide.shapes.add_shape(1, Inches(6.5), Inches(5), Inches(5.8), Inches(1.5))
hero_box.fill.solid()
hero_box.fill.solid()
hero_box.fill.fore_color.rgb = ACCENT_COLOR
hero_box.line.fill.background()
tf = hero_box.text_frame
p = tf.paragraphs[0]
p.text = "The Hero\nTaste Infrastructure\n\"We can prove it's safe.\""
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.bold = True


# --------------------------------------------------------------------------------
# SLIDE 3: INVERSION (Dark) - Graph Slide
# --------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)

add_text_box(slide, "02 — The Inversion", Inches(1), Inches(0.5), Inches(5), Inches(0.5), 12, TEXT_GRAY)
add_text_box(slide, "Stated precisely", Inches(1), Inches(1), Inches(10), Inches(1), 36, TEXT_DARK, bold=True)

# Graph container area
# Draw axes
connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(2.5), Inches(1.5), Inches(6))
connector.line.color.rgb = TEXT_GRAY
connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(6), Inches(10), Inches(6))
connector.line.color.rgb = TEXT_GRAY

# Draw "Generation" line (Going down)
line_gen = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(1.5), Inches(3), Inches(10), Inches(5.8))
line_gen.line.color.rgb = TEXT_GRAY
line_gen.line.dash_style = MSO_LINE_DASH_STYLE.DASH
line_gen.line.width = Pt(2)
add_text_box(slide, "Generation Cost (Collapsing)", Inches(10.1), Inches(5.6), Inches(3), Inches(0.5), 12, TEXT_GRAY)

# Draw "Evaluation" line (Going up/flat)
line_eval = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(1.5), Inches(3.2), Inches(10), Inches(3))
line_eval.line.color.rgb = ACCENT_COLOR
line_eval.line.width = Pt(4)
add_text_box(slide, "Evaluation Cost (Sticky)", Inches(10.1), Inches(2.8), Inches(3), Inches(0.5), 12, ACCENT_COLOR, bold=True)

# Danger Zone annotation
add_text_box(slide, "THE DANGER ZONE:\nWhere speed kills quality", Inches(5), Inches(3.5), Inches(3), Inches(1), 12, ACCENT_COLOR, align=PP_ALIGN.CENTER)

# List points
add_text_box(slide, "1. Marginal cost of generating is collapsing.\n2. Marginal cost of validating is NOT collapsing.", 
             Inches(1), Inches(6.5), Inches(11), Inches(1), 16, TEXT_DARK)


# --------------------------------------------------------------------------------
# SLIDE 4: JUDGMENT (Light)
# --------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)

add_text_box(slide, "03 — What Judgment Really Is", Inches(1), Inches(0.5), Inches(5), Inches(0.5), 12, TEXT_GRAY)
add_text_box(slide, "The Luxury Parallel: Hermès doesn't ship \"mostly good\" bags.", Inches(1), Inches(1), Inches(11), Inches(1.2), 32, TEXT_LIGHT, bold=True)

add_text_box(slide, "In AI, we can BUILD judgment systems. Scale them. Prove them.", Inches(1), Inches(2.5), Inches(10), Inches(1), 20, ACCENT_COLOR, bold=True)

# Comparison Boxes
# Box 1
box1 = slide.shapes.add_shape(1, Inches(1), Inches(4), Inches(5), Inches(2))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(240, 240, 240)
box1.line.fill.background()
tf = box1.text_frame
tf.text = "Training time\n5-20 years\n(Human)"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[1].font.size = Pt(24)
tf.paragraphs[1].font.bold = True
tf.paragraphs[1].font.color.rgb = TEXT_LIGHT

# Box 2
box2 = slide.shapes.add_shape(1, Inches(6.5), Inches(4), Inches(5), Inches(2))
box2.fill.solid()
box2.fill.fore_color.rgb = ACCENT_COLOR
box2.line.fill.background()
tf = box2.text_frame
tf.text = "Inference time\nMilliseconds\n(System)"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = TEXT_DARK
tf.paragraphs[1].font.size = Pt(24)
tf.paragraphs[1].font.bold = True
tf.paragraphs[1].font.color.rgb = TEXT_DARK
tf.paragraphs[2].font.color.rgb = TEXT_DARK

add_text_box(slide, "\"Judgment can be instant, but it is rarely free. Someone paid for it in years.\"", 
             Inches(1), Inches(6.2), Inches(11), Inches(1), 16, TEXT_GRAY) # Italic


# --------------------------------------------------------------------------------
# SLIDE 5: ECONOMICS (Dark)
# --------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)

add_text_box(slide, "04 — Economics", Inches(1), Inches(0.5), Inches(5), Inches(0.5), 12, TEXT_GRAY)
add_text_box(slide, "Why value accrues to judgment now", Inches(1), Inches(1), Inches(10), Inches(1), 36, TEXT_DARK, bold=True)

add_text_box(slide, "In markets, rents concentrate at constraints.", Inches(1), Inches(2.2), Inches(10), Inches(0.5), 18, TEXT_DARK)

# Box 1: $10 Rule
box1 = slide.shapes.add_shape(1, Inches(1), Inches(3.5), Inches(5), Inches(2.5))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(50, 50, 50)
box1.line.color.rgb = RGBColor(100, 100, 100)
tf = box1.text_frame
tf.text = "The $10 Rule\n\nEvery dollar you save on cheap generation, you will spend $10 on cleaning up the mess if you don't have evaluation gates."
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = ACCENT_COLOR

# Box 2: Shifting Constraints
box2 = slide.shapes.add_shape(1, Inches(6.5), Inches(3.5), Inches(5), Inches(2.5))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(50, 50, 50)
box2.line.color.rgb = RGBColor(100, 100, 100)
tf = box2.text_frame
tf.text = "Shifting Constraints\n\n2010: Distribution (Netflix)\n2020: Content (Creators)\n2025: Selection (Curators win)"
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = ACCENT_COLOR


# --------------------------------------------------------------------------------
# SLIDE 6: BUSINESS CASE (Light)
# --------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)

add_text_box(slide, "05 — The Enterprise Case", Inches(1), Inches(0.5), Inches(5), Inches(0.5), 12, TEXT_GRAY)
add_text_box(slide, "The Cost of Ignorance", Inches(1), Inches(1), Inches(10), Inches(1), 36, TEXT_LIGHT, bold=True)

# Nightmare Scenario
add_text_box(slide, "The Nightmare Scenario", Inches(1), Inches(2.5), Inches(5), Inches(0.5), 16, ACCENT_COLOR, bold=True)
add_text_box(slide, "Your competitor just shipped 10,000 AI descriptions. Half are wrong. Revenue drops 15%. They can't trace why.", 
             Inches(1), Inches(3.2), Inches(5), Inches(2), 14, TEXT_LIGHT)

# ROI Math
add_text_box(slide, "THE SIMPLE MATH (ROI)", Inches(6.5), Inches(2.5), Inches(5), Inches(0.5), 16, ACCENT_COLOR, bold=True)
add_text_box(slide, "• Cost to build gate: $200k-500k\n• Cost of brand disaster: $5M-50M\n• Break-even probability: 1%", 
             Inches(6.5), Inches(3.2), Inches(5), Inches(2), 14, TEXT_LIGHT)


# --------------------------------------------------------------------------------
# SLIDE 7: ARBITRAGE (Light) - Graph Slide
# --------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)

add_text_box(slide, "06 — The Arbitrage", Inches(1), Inches(0.5), Inches(5), Inches(0.5), 12, TEXT_GRAY)
add_text_box(slide, "Where the market is still asleep", Inches(1), Inches(1), Inches(10), Inches(1), 36, TEXT_LIGHT, bold=True)

# Graph axes
connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(2.5), Inches(1.5), Inches(6))
connector.line.color.rgb = TEXT_GRAY
connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(6), Inches(10), Inches(6))
connector.line.color.rgb = TEXT_GRAY

# Tech Skills Line (Going Down)
line_tech = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(2.5), Inches(10), Inches(5.8))
line_tech.line.color.rgb = TEXT_GRAY
line_tech.line.dash_style = MSO_LINE_DASH_STYLE.DASH
add_text_box(slide, "Tech Skills Value", Inches(10.1), Inches(5.6), Inches(2), Inches(0.5), 10, TEXT_GRAY)

# Taste/Judgment Line (Going Up)
line_taste = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.5), Inches(5.8), Inches(10), Inches(2.5))
line_taste.line.color.rgb = ACCENT_COLOR
line_taste.line.width = Pt(4)
add_text_box(slide, "Judgment Value", Inches(10.1), Inches(2.3), Inches(2), Inches(0.5), 10, ACCENT_COLOR, bold=True)

# Crossover annotation
shape = slide.shapes.add_shape(9, Inches(5.5), Inches(3.8), Inches(0.4), Inches(0.4)) # Oval
shape.fill.background()
shape.line.color.rgb = ACCENT_COLOR
add_text_box(slide, "CROSSOVER: NOW", Inches(5), Inches(3.2), Inches(2), Inches(0.5), 12, ACCENT_COLOR, bold=True)

# Boxes at bottom
add_text_box(slide, "CTO: Hire for taste.", Inches(1), Inches(6.5), Inches(3), Inches(0.8), 12, TEXT_LIGHT, bold=True)
add_text_box(slide, "INVESTOR: Bet on Eval Infra.", Inches(5), Inches(6.5), Inches(3), Inches(0.8), 12, TEXT_LIGHT, bold=True)
add_text_box(slide, "FOUNDER: Hire judges.", Inches(9), Inches(6.5), Inches(3), Inches(0.8), 12, TEXT_LIGHT, bold=True)


# --------------------------------------------------------------------------------
# SLIDE 8: STACK (Dark)
# --------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)

add_text_box(slide, "07 — The Gatekeeper Model", Inches(1), Inches(0.5), Inches(5), Inches(0.5), 12, TEXT_GRAY)
add_text_box(slide, "Your stack is a funnel, not a checklist.", Inches(1), Inches(1), Inches(10), Inches(1), 30, TEXT_DARK, bold=True)

# Code/Funnel Block
box = slide.shapes.add_shape(1, Inches(1), Inches(2.2), Inches(11.3), Inches(4))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(40, 40, 40)
box.line.color.rgb = RGBColor(60, 60, 60)
tf = box.text_frame
content = """INPUT: 1,000 Candidates
↓
GATE 1: SAFETY (PII, Regex) -> 99.8% Survival
↓
GATE 2: TECHNICAL (Hallucinations) -> 94% Survival
↓
GATE 3: BRAND FIT (LLM Judge) -> 87% Survival
↓
GATE 4: EXCELLENCE (Reward Model) -> 78% Survival
↓
SHIP: 64% Final Pass Rate"""
tf.text = content
for p in tf.paragraphs:
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(200, 200, 200)

add_text_box(slide, "The Litmus Test: If you cannot automatically reject the bottom 30%, you don't have infrastructure. You have hope.", 
             Inches(1), Inches(6.4), Inches(11), Inches(0.8), 14, ACCENT_COLOR)


# --------------------------------------------------------------------------------
# SLIDE 9: PREDICTIONS (Light)
# --------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)

add_text_box(slide, "08 — Falsifiable Bets", Inches(1), Inches(0.5), Inches(5), Inches(0.5), 12, TEXT_GRAY)
add_text_box(slide, "Career-Defining Stakes", Inches(1), Inches(1), Inches(10), Inches(1), 36, TEXT_LIGHT, bold=True)

# Table
rows = 4
cols = 2
left = Inches(1)
top = Inches(2.5)
width = Inches(11.3)
height = Inches(3)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(3)
table.columns[1].width = Inches(8.3)

# Headers
table.cell(0, 0).text = "The Bet"
table.cell(0, 1).text = "The Stakes"

# Row 1
table.cell(1, 0).text = "The Talent Filter"
table.cell(1, 1).text = "By 2027, no F500 will promote an AI lead who cannot explain their eval metrics."

# Row 2
table.cell(2, 0).text = "The Budget Flip"
table.cell(2, 1).text = "By 2028, enterprises will spend 3x more on Evaluation than on Foundation Models."

# Row 3
table.cell(3, 0).text = "The New Unicorn"
table.cell(3, 1).text = "The next $10B AI company won't be a model lab. It will be an Eval Infra company."

# Style Table Text
for row in table.rows:
    for cell in row.cells:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = TEXT_LIGHT


# --------------------------------------------------------------------------------
# SLIDE 10: CREDIBILITY (Dark)
# --------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)

add_text_box(slide, "09 — Skin in the Game", Inches(1), Inches(0.5), Inches(5), Inches(0.5), 12, TEXT_GRAY)
add_text_box(slide, "Why this is my life's work.", Inches(1), Inches(1), Inches(10), Inches(1), 36, TEXT_DARK, bold=True)

add_text_box(slide, "Before my current role: 16 years building taste systems for luxury brands (Fendi, Bentley). I'm building W230 to fix the AI quality gap.", 
             Inches(1), Inches(2), Inches(11), Inches(1), 16, RGBColor(200, 200, 200))

# Left Box: The Save
boxL = slide.shapes.add_shape(1, Inches(1), Inches(3.2), Inches(5.5), Inches(3))
boxL.fill.background()
boxL.line.color.rgb = ACCENT_COLOR
tf = boxL.text_frame
tf.text = "THE SAVE (CONCRETE PROOF)\n\nFinding: 8% factual errors in baseline.\nExposure: $15M potential risk.\nResult: 125x ROI. Quality became the edge."
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = ACCENT_COLOR
tf.paragraphs[1].font.color.rgb = TEXT_DARK
tf.paragraphs[2].font.color.rgb = TEXT_DARK
tf.paragraphs[3].font.color.rgb = TEXT_DARK

# Right Box: Background
boxR = slide.shapes.add_shape(1, Inches(6.8), Inches(3.2), Inches(5.5), Inches(3))
boxR.fill.background()
boxR.line.color.rgb = TEXT_GRAY
tf = boxR.text_frame
tf.text = "BACKGROUND\n\nLuxury Brands (Fendi, Bentley, LVMH)\nStudio W230 (Since 2016)\nMIT Applied Agentic AI (2025)"
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT_GRAY
for p in tf.paragraphs[1:]:
    p.font.color.rgb = TEXT_DARK


# --------------------------------------------------------------------------------
# SLIDE 11: CLAIM (Dark)
# --------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)

add_text_box(slide, "10 — The Claim", Inches(1), Inches(0.5), Inches(5), Inches(0.5), 12, TEXT_GRAY)

# Big centered Text
center_box = add_text_box(slide, "You are either building this now,\nor you are cleaning up later.", 
             Inches(1.5), Inches(2), Inches(10.3), Inches(2), 32, TEXT_DARK, align=PP_ALIGN.CENTER)

# "CHOOSE"
choose_box = add_text_box(slide, "Choose.", 
             Inches(1.5), Inches(3.8), Inches(10.3), Inches(1.5), 80, ACCENT_COLOR, bold=True, align=PP_ALIGN.CENTER)

# Contact
add_text_box(slide, "Studio W230 — Invite-only evaluation audits launching 2026.\nRegister: simone@w230.net", 
             Inches(1.5), Inches(6), Inches(10.3), Inches(1), 16, TEXT_GRAY, align=PP_ALIGN.CENTER)


# --------------------------------------------------------------------------------
# SLIDE 12: REFERENCES (Light)
# --------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)

add_text_box(slide, "Further Reading", Inches(1), Inches(0.5), Inches(5), Inches(0.5), 12, TEXT_GRAY)
add_text_box(slide, "For those who want to go deeper.", Inches(1), Inches(1), Inches(10), Inches(1), 32, TEXT_LIGHT, bold=True)

content = """
Kahneman (2011) — Thinking, Fast and Slow
Simon (1971) — Designing Organizations for an Information-Rich World
Gartner (2025) — The Case for AI Quality Gates
Anthropic — Model Evaluation & Constitutional AI
"""
add_text_box(slide, content, Inches(1), Inches(2.5), Inches(10), Inches(4), 18, TEXT_LIGHT)


# --------------------------------------------------------------------------------
# SAVE
# --------------------------------------------------------------------------------
prs.save('Taste_Infrastructure_W230.pptx')
print("Presentation generated successfully: Taste_Infrastructure_W230.pptx")